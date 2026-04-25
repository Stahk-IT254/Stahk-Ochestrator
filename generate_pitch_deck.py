import os
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Please install python-pptx first by running:")
    print("pip install python-pptx")
    exit()

def add_slide(prs, title_text, content_bullets):
    slide_layout = prs.slide_layouts[1] # 1 is Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    # Format Title
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.bold = True
        
    content = slide.placeholders[1]
    tf = content.text_frame
    
    for i, bullet in enumerate(content_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        p.space_after = Pt(14)
        
    return slide

def create_pitch_deck():
    prs = Presentation()

    # SLIDE 1: TITLE
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Stahk Orchestrator"
    subtitle.text = "The Decentralized Marketplace for the Agentic Economy\nBuilt on ARC & Circle USDC"

    # SLIDE 2: THE PROBLEM
    add_slide(prs, "The Problem", [
        "Rigid Subscriptions: Users pay expensive monthly SaaS fees for AI tools they underutilize.",
        "Monetization Gap: Independent AI developers lack infrastructure to monetize specialized agents.",
        "Hallucinations: Generic LLMs lack domain-specific truth, making them unreliable for serious tasks.",
        "Friction: Coordinating multiple AI agents across different APIs is complex and manual."
    ])

    # SLIDE 3: THE SOLUTION
    add_slide(prs, "The Solution: Stahk Orchestrator", [
        "Agentic Marketplace: Users submit real-world goals and hire specialized, autonomous AI agents.",
        "Micro-Transactions: Users pay fractions of a cent (USDC) only for the compute they use.",
        "Trustless Escrow: Funds are held safely and only released when the agent delivers the result.",
        "Domain-Aware RAG: Agents query specific knowledge bases (Finance, Construction) to guarantee factual outputs."
    ])

    # SLIDE 4: ARCHITECTURE & TECH STACK
    add_slide(prs, "Architecture & Tech Stack", [
        "Frontend: Next.js + React (Dynamic UI, Auto-polling)",
        "Backend: Django + Python (Orchestration Engine)",
        "Financial Layer: Circle Programmable Wallets + USDC (Nanopayments)",
        "Intelligence Layer: ChromaDB RAG Pipeline",
        "Quality Control: Automated API latency & integrity vetting"
    ])

    # SLIDE 5: HOW IT WORKS (USER FLOW)
    add_slide(prs, "How It Works", [
        "1. Submit a Goal: User inputs a task (e.g., 'Plan a residential construction').",
        "2. Hire Agents: Orchestrator matches user with top-rated domain agents.",
        "3. Escrow Lock: USDC nanopayment is locked in escrow.",
        "4. Autonomous Execution: Orchestrator fires agents in the background.",
        "5. Payout & PDF Export: Platform takes a 10% matching fee, creator is paid, user exports the final report."
    ])

    # SLIDE 6: BUSINESS MODEL
    add_slide(prs, "Business Model & Scalability", [
        "10% Matching Fee: Orchestrator automatically deducts a 10% fee on every transaction.",
        "Infinite Supply: Any developer can connect their Agent API, pass the automated vetting, and start earning.",
        "Nanopayment Advantage: High-frequency, sub-cent pricing unlocks a completely new AI economy."
    ])

    # SLIDE 7: THANK YOU
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You!"
    subtitle.text = "Welcome to the Agentic Economy."

    output_path = "Stahk_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Presentation generated successfully: {output_path}")

if __name__ == "__main__":
    create_pitch_deck()
